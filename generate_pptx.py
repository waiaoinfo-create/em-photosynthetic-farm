#!/usr/bin/env python3
"""
Generate EM菌 + 光合菌 Farm Practical Guide PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
C_GREEN_DARK = RGBColor(0x1B, 0x43, 0x32)
C_GREEN_MID = RGBColor(0x2D, 0x6A, 0x4F)
C_GREEN_LIGHT = RGBColor(0x40, 0x91, 0x6C)
C_BG_LIGHT = RGBColor(0xF0, 0xF7, 0xF4)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT = RGBColor(0x95, 0xD5, 0xB2)
C_DANGER = RGBColor(0xE7, 0x6F, 0x51)
C_WARNING = RGBColor(0xE9, 0xC4, 0x6A)
C_ORANGE = RGBColor(0xF4, 0xA2, 0x61)
C_TEXT = RGBColor(0x1B, 0x43, 0x32)

def add_bg(slide, color=C_GREEN_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val=str(int(alpha * 1000)))
    return shape

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, C_GREEN_DARK)

    # Decorative bar
    add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), C_ACCENT)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EM菌與光合菌"
    p.font.size = Pt(48)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "農場實務操作完全手冊"
    p2.font.size = Pt(36)
    p2.font.color.rgb = C_ACCENT
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\n從基礎理論到實務操作 · 成本分析 · 人力配置 · 真實案例 · 極致拆解"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0xB7, 0xE4, 0xC7)
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "\n2026年6月  |  資料來源：食工所、農業試驗所、農業知識入口網、上下游新聞"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0x88, 0xAA, 0x99)
    p4.alignment = PP_ALIGN.CENTER

def add_section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_GREEN_MID)
    add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.04), C_ACCENT)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"第{number}部分"
    p.font.size = Pt(20)
    p.font.color.rgb = C_ACCENT
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.color.rgb = C_WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

    if subtitle:
        p3 = tf2.add_paragraph()
        p3.text = subtitle
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(0xB7, 0xE4, 0xC7)
        p3.alignment = PP_ALIGN.CENTER

def add_content_slide(title, left_items, right_items=None, bullet_color=C_GREEN_MID):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_LIGHT)

    # Title bar
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_GREEN_MID)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = C_WHITE
    p.font.bold = True

    def add_items(left, top, width, items):
        txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            if isinstance(item, tuple):
                # (bold_part, normal_part)
                run1 = p.add_run()
                run1.text = item[0]
                run1.font.size = Pt(16)
                run1.font.bold = True
                run1.font.color.rgb = C_GREEN_MID
                run2 = p.add_run()
                run2.text = item[1]
                run2.font.size = Pt(16)
                run2.font.color.rgb = C_TEXT
            else:
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = C_TEXT
            p.space_after = Pt(8)
            p.level = 0

    if right_items:
        # Two columns
        add_items(Inches(0.5), Inches(1.2), Inches(5.8), left_items)
        # Divider
        add_shape_bg(slide, Inches(6.5), Inches(1.2), Inches(0.03), Inches(5.5), C_ACCENT)
        add_items(Inches(6.8), Inches(1.2), Inches(5.8), right_items)
    else:
        add_items(Inches(0.5), Inches(1.2), Inches(12.3), left_items)

    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_BG_LIGHT)

    # Title bar
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_GREEN_MID)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = C_WHITE
    p.font.bold = True

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(12.3)
    height = Inches(0.4 * n_rows + 0.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = C_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_GREEN_MID

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = C_TEXT
                paragraph.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if ri % 2 == 0 else RGBColor(0xE8, 0xF5, 0xE9)

    return slide

def add_highlight_slide(title, items, bg_color=C_GREEN_MID):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)

    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.04), C_ACCENT)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        if isinstance(item, tuple):
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(20)
            run1.font.bold = True
            run1.font.color.rgb = C_ACCENT
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(20)
            run2.font.color.rgb = C_WHITE
        else:
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = C_WHITE
        p.space_after = Pt(12)
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
add_title_slide()

# ============================================================
# SLIDE 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_BG_LIGHT)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), C_GREEN_MID)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "目錄"
p.font.size = Pt(28)
p.font.color.rgb = C_WHITE
p.font.bold = True

toc_items = [
    "1.  EM菌與光合菌概論",
    "2.  核心差異釐清",
    "3.  光合菌培養SOP",
    "4.  EM菌自製SOP",
    "5.  農場四大應用領域",
    "6.  成本量化分析",
    "7.  人力配置實務",
    "8.  真實案例深度拆解",
    "9.  常見失敗與排除",
    "10. 風險與注意事項",
    "11. 總結效益矩陣",
]
txBox2 = slide.shapes.add_textbox(Inches(3.5), Inches(1.3), Inches(6.3), Inches(5.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(toc_items):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = C_GREEN_MID
    p.space_after = Pt(10)
    p.font.bold = True

# ============================================================
# SLIDE 3-4: Section 1 - 概論
# ============================================================
add_section_slide("壹", "EM菌與光合菌概論")

add_content_slide("EM菌 vs 光合菌 基礎認知", [
    ("EM菌（Effective Microorganisms）\n", "日本琉球大學比嘉照夫教授研發\n5大類10屬56餘種微生物複合\n光合菌+乳酸菌+酵母菌+放線菌+絲狀菌\n綜合型：土壤改良、飼料添加、環境除臭"),
    ("光合菌（PSB）\n", "紫色非硫菌——沼澤紅假單胞菌為主\n固氮溶磷、分解氨氮與硫化氫\n分泌植物生長因子、促進根系發育\n台灣食工所（FIRDI）主要推廣"),
])

# ============================================================
# SLIDE 5: 核心差異表
# ============================================================
add_table_slide("核心差異釐清", ["比較項目", "EM菌", "光合菌"], [
    ["取得難度", "容易（市售多）", "需特定來源"],
    ["自製難度", "低（糖+水+菌種）", "中（需光照管理）"],
    ["土壤改良", "全面性佳", "專一性強（固氮溶磷）"],
    ["水產應用", "良好", "極佳（降氨氮）"],
    ["畜禽除臭", "極佳", "良好"],
    ["葉面噴施", "佳", "極佳（促根系）"],
    ["堆肥加速", "極佳", "普通"],
    ["保存性", "較差（液體衰退快）", "較佳（厭氧穩定）"],
])

# ============================================================
# SLIDE 6-8: 光合菌SOP
# ============================================================
add_section_slide("參", "光合菌培養SOP", "以食工所（FIRDI）技術為準")

add_table_slide("光合菌材料清單（20公升）", ["項目", "數量", "單價", "小計"], [
    ["光合菌母菌", "200ml", "0~500元", "0~500元"],
    ["RO水/蒸餾水", "20L桶", "20~35元", "20~35元"],
    ["培養基粉末", "200g包", "90~150元", "90~150元"],
    ["透明PE桶", "1個", "50~80元", "50~80元*"],
    ["每桶總成本", "", "", "110~185元"],
    ["每公升成本", "", "", "5.5~9.3元"],
])

add_content_slide("光合菌培養步驟", [
    ("第一步：準備\n", "20L透明桶 → 75%酒精消毒 → RO水或煮沸除氯自來水 → 水溫25~35°C"),
    ("第二步：混合\n", "先加5L水 → 加培養基200g搖晃溶解 → 加母菌200ml → 加滿至桶口（少留空氣）→ 旋緊"),
    ("第三步：培養\n", "陽光充足處（避免>40°C）→ 每日搖晃1~2次 → 25~35°C → 觀察顏色：淡粉→粉紅→紫紅\n夏季7~10天，冬季14~21天"),
    ("第四步：完成判別\n", "深紅/紫紅色，醬油香或烏梅香（非臭味）→ 陰涼處密封保存3~6個月"),
])

add_highlight_slide("光合菌品質判別", [
    ("✅ 成功：", "暗紅/紫紅色 · 醬油香/烏梅香 · pH 5.5~6.5 · 少許沉澱"),
    ("❌ 失敗：", "灰黑/綠色混濁 · 惡臭/霉味 · pH<4.5或>8.0 · 大量不明沉澱"),
])

# ============================================================
# SLIDE 9-10: EM菌SOP
# ============================================================
add_section_slide("肆", "EM菌自製SOP")

add_table_slide("EM菌自產法（40公斤產量）", ["材料", "數量", "成本"], [
    ["紅糖", "2~4kg", "80~160元"],
    ["水（除氯）", "40L", "0元"],
    ["EM活菌自產劑", "500g", "150~250元"],
    ["50L塑膠桶（首次）", "1個", "200~300元"],
    ["每公斤成本（不含桶）", "", "5.8~10.3元"],
])

add_content_slide("EM菌操作步驟與紅糖擴繁法", [
    ("EM活菌自產法\n", "① 桶中加35°C溫水40L\n② 加紅糖2~4kg攪拌溶解\n③ 加EM自產劑500g攪勻\n④ 密封25~35°C培養\n⑤ 每日開蓋放氣+攪拌\n⑥ 約7天完成（產氣連續減少）"),
    ("紅糖擴繁法（簡易版）\n", "市售EM原液1L + 紅糖1kg + 水20L → 密封5~7天\n⚠ 菌數逐代衰減，最多擴繁3代即需重購"),
])

# ============================================================
# SLIDE 11-15: 四大應用
# ============================================================
add_section_slide("伍", "農場四大應用領域", "種植 · 水產 · 畜禽 · 堆肥")

add_table_slide("種植業應用——土壤改良與葉面噴施", ["項目", "EM菌", "光合菌"], [
    ["整地前土壤", "稀釋500倍，200~300L/分地", "稀釋300倍，150~200L/分地"],
    ["浸種", "稀釋100~200倍，15~30分鐘", "稀釋100~200倍，15~30分鐘"],
    ["定植澆根", "200~300倍，200~500ml/株", "300~500倍，300~500ml/株"],
    ["葉菜類噴施", "300~500倍，7~10天/次", "300~500倍，7~10天/次"],
    ["果樹噴施", "200~300倍，10~14天/次", "200~300倍，10~14天/次（避開花期）"],
    ["根莖類", "500~800倍，僅幼苗期", "500~800倍，膨大期禁用"],
])

add_content_slide("水產與畜禽養殖應用", [
    ("水產養殖\n", "EM菌：全池1~2ppm/周 · 拌料5~10ml/kg\n光合菌：全池5~10ppm/日 · 拌料10~20ml/kg\n光合菌對氨氮（NH₃-N）分解效果優於EM菌"),
    ("畜禽養殖\n", "飲水：500~1000倍稀釋\n飼料發酵：EM菌2~5L/噸+糖2kg+水100L\n噴霧除臭：稀釋100~200倍，每周2~3次\n糞便堆肥：EM菌5~10L/噸+米糠20kg，15~30天\n料肉比↓15%，氨氣↓60~80%，死亡率↓30~50%"),
])

add_content_slide("堆肥製作與生物防蟲液", [
    ("EM菌堆肥\n", "有機物料1噸 + EM菌5~10L + 紅糖5kg\n水分50~60% → 堆置1~1.5m高 → 覆蓋帆布\n7天後翻堆，之後每3~5天翻一次\n夏季20~30天／冬季40~60天完全腐熟\n對比傳統：時間↓50~60%，臭味↓80~90%"),
    ("生物防蟲液配方\n", "EM菌1L + 紅糖1L + 醋精1L + 白酒1L + 辣椒粉1L + 水10L\n密封發酵15天（25°C）\n使用時稀釋200~1000倍噴灑葉面\n對蚜蟲、青蟲、部分真菌病害有防治效果"),
])

# ============================================================
# SLIDE 16-18: 成本分析
# ============================================================
add_section_slide("陸", "成本量化分析", "自製 vs 市售 · ROI > 500%")

add_table_slide("自製 vs 市售價格比較", ["項目", "市售價格", "自製價格", "節省比例"], [
    ["EM菌液（1L）", "15~30元", "1~3元", "90~93%"],
    ["光合菌液（1L）", "100~130元", "5~10元", "90~95%"],
    ["20L光合菌", "2,500元", "110~185元", "92~96%"],
])

add_content_slide("一分地年度成本 vs 真實案例", [
    ("一分地（0.1公頃）年度菌液成本\n", "EM菌年用量100~200L\n自製：110~660元  vs  購買：1,500~6,000元\n\n光合菌年用量100~200L\n自製：550~1,850元  vs  購買：10,000~26,000元"),
    ("真實案例：台南紅龍果——黃政德（4分地）\n", "使用光合菌前：年肥料費 60,000元\n使用光合菌後：年肥料費 12,000元\n每年節省 48,000元（每分地省12,000元）\n\n首次投入：700~1,500元\n年度經常成本：4,000~8,000元\n年度節省效益：40,000~80,000元"),
])

add_highlight_slide("投資回報率 ROI >> 500% / 年", [
    ("首次投入：", "700~1,500元（母菌+培養基+容器）"),
    ("年度經常成本：", "4,000~8,000元（培養基補充+糖+水）"),
    ("年度節省效益：", "肥料減少30,000~50,000元\n+ 農藥減少5,000~15,000元\n+ 產量提升10~30%"),
    ("結論：", "每投入1元，回報5元以上"),
])

# ============================================================
# SLIDE 19-21: 人力配置
# ============================================================
add_section_slide("柒", "人力配置實務", "3~5公頃農場 · 月增40~60工時")

add_table_slide("常態工作量（3~5公頃農場）", ["工作項目", "頻率", "每次工時", "月工時"], [
    ["菌液培養", "月1~2次", "1~2小時", "2~4小時"],
    ["品質檢查", "周1次", "0.5小時", "2小時"],
    ["稀釋調配", "每次使用前", "0.5小時", "2~4小時"],
    ["葉面噴施", "周1次", "2~4小時/公頃", "8~16小時"],
    ["灌根澆灌", "每2周1次", "1~2小時/公頃", "2~4小時"],
    ["堆肥翻堆", "月1~2次", "2~4小時", "4~8小時"],
    ["畜舍噴霧", "周2次", "0.5~1小時", "4~8小時"],
    ["記錄監測", "每日", "0.5小時", "15小時"],
    ["合計", "", "", "40~60小時"],
])

add_content_slide("導入期人力配置與設備需求", [
    ("導入期前3個月\n", "第1~2周：設備採購、場地規劃（場長+1助手 40小時）\n第1個月：首次菌液培養、小面積試驗（25小時）\n第2~3個月：逐步擴大至30%面積（40小時）\n第4~6個月：全面導入，回歸常態"),
    ("設備需求\n", "HDPE桶 50~200L：200~800元/個\n20L透明桶：50~100元/個\n量筒/量杯：100~200元\npH計/試紙：200~1,000元\n動力噴霧機：5,000~20,000元\n背負式噴霧器：1,000~3,000元\n顯微鏡選配：3,000~8,000元"),
])

# ============================================================
# SLIDE 22-25: 真實案例
# ============================================================
add_section_slide("捌", "真實案例深度拆解", "台灣農民第一手實證數據")

add_content_slide("案例一：台南善化紅龍果——黃政德", [
    ("背景\n", "4分地紅龍果園，2016年導入光合菌\n原慣行農法全面轉型\n試驗兩行對照：光合菌無肥 vs 慣行施肥"),
    ("成效數據\n", "實驗組單果重：2公斤（對照組0.8~1.2kg）\n產期延長至冬季仍可產出\n肥料用量降至原有的 1/5\n年肥料費：60,000元 → 12,000元\n材料成本：第一年25,000元 → 第三年後3,000元/年"),
    ("教訓\n", "根莖類作物膨大期禁止使用\n開花前最佳使用時機\n長期需搭配基礎施肥"),
])

add_content_slide("案例二：屏東芒果——林毓健", [
    ("背景\n", "2018年參加食工所光合菌課程\n以往市售光合菌 20L = 2,500元\n自製成本不到 200元/20L"),
    ("操作\n", "稀釋300~500倍，每月2次灌根\n避開開花期（否則長葉不開花）\n不混用殺菌劑"),
    ("成效\n", "芒果樹勢健康、葉片旺盛\n連續使用3年效果累積\n「市售光合菌很混亂，不敢亂買」"),
])

add_content_slide("案例三：中國EM菌養豬場 + 案例四：屏東縣政府", [
    ("EM菌養豬場實證數據\n", "料肉比 ↓ 15%\n畜舍氨氣濃度 ↓ 60~80%\n死亡率 ↓ 30~50%\n獸藥費用 ↓ 40~60%\n\n水產：每畝每周潑灑1kg菌液\n飼料添加：100~200ml + 糖200g + 水15kg"),
    ("屏東縣政府推廣光合菌\n", "縣府農業處邀請食工所廖麗玲博士指導\n光合菌具固氮溶磷作用\n增加土壤益菌量\n提高作物吸收力與品質\n達農藥及肥料減量目標"),
])

# ============================================================
# SLIDE 26-27: 常見失敗
# ============================================================
add_section_slide("玖", "常見失敗原因與排除對策")

add_table_slide("培養失敗與田間應用失敗對策", ["現象", "原因", "解決方案"], [
    ["菌液發黑/發臭", "雜菌污染", "重新培養，加強消毒，使用RO水"],
    ["菌液不變紅", "光照不足/溫度太低", "移至陽光充足處，保持25°C以上"],
    ["菌液變綠", "藻類污染", "加強密封避光，檢討水源"],
    ["無明顯效果", "濃度不足/頻率太低", "提高至200~300倍，縮短至7天/次"],
    ["病害增加", "混入病原菌", "送檢菌液，加強無菌操作"],
    ["葉片灼傷", "稀釋<100倍/中午噴", "300倍以上，傍晚噴施"],
    ["花果掉落", "開花期濃度過高", "暫停或稀釋至800倍以上"],
])

# ============================================================
# SLIDE 28-29: 風險
# ============================================================
add_section_slide("拾", "風險與注意事項")

add_content_slide("8大注意事項", [
    ("① 光合菌不是肥料\n", "它分解養分而非創造養分，需搭配基礎施肥"),
    ("② 菌液保存期限\n", "EM菌1個月內最佳（2個月後活菌降50%）\n光合菌可保存3~6個月"),
    ("③ 殺菌劑禁忌\n", "不可與殺菌劑、抗生素混合使用"),
    ("④ 法規現況\n", "光合菌尚未正式登記為微生物肥料\n（非孢子菌保存不易，無廠商登記）\n市售產品品質參差不齊"),
    ("⑤ 根莖類限制\n", "地下部膨大期禁用光合菌"),
    ("⑥ 開花期限制\n", "光合菌盛花期暫停使用"),
    ("⑦ 雜菌污染\n", "自製菌液可能混入病原菌，加強消毒"),
    ("⑧ 購買建議\n", "向食工所等可信來源取得母菌"),
])

# ============================================================
# SLIDE 30: 總結
# ============================================================
add_highlight_slide("總結：最終建議", [
    ("年淨效益：", "8~15萬元（3~5公頃農場）"),
    ("首次投入：", "低於2,000元"),
    ("月增人力：", "40~60工時（分散於現有勞動力）"),
    ("最高策略：", "光合菌為主（根系強化）+ EM菌為輔（堆肥+除臭）"),
    ("入門行動：", "先向食品工業發展研究所參加免費課程，取得正確母菌與技術"),
])

add_highlight_slide("資料來源", [
    ("財團法人食品工業發展研究所（FIRDI）", ""),
    ("農業部農業試驗所", ""),
    ("農業部農業知識入口網（kmweb.moa.gov.tw）", ""),
    ("上下游新聞（www.newsmarket.com.tw）", "林怡均記者報導"),
    ("水產試驗所", ""),
    ("農業用微生物產業固本與加值應用技術研發成果", "農業部2021-2024統籌計畫"),
], bg_color=C_GREEN_DARK)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "EM菌與光合菌農場實務手冊.pptx")
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
